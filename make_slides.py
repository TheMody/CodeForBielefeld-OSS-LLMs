# -*- coding: utf-8 -*-
"""
Build "Open Source LLMs" deck for Code for Bielefeld from plan.md.
Generates OpenSource-LLMs-CodeForBielefeld.pptx and (optionally) PNG previews
of every slide for QA (--preview <dir>).

Run:  py -3.13 make_slides.py [--preview <dir>]
"""
import os, sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

from PIL import Image, ImageDraw, ImageFont

# ---------------------------------------------------------------- constants
EMU_IN = 914400
W, H = 13.333, 7.5
SCALE = 130  # px per inch for previews

BG      = "20232D"
CARD    = "2A2E3C"
CARD2   = "262A37"
TERM    = "14161E"
TERMBRD = "3A3F52"
CORAL   = "FF867F"
CORALDK = "8A4A46"
TEXT    = "EEF0F6"
MUTED   = "9AA3B8"
FAINT   = "6F7differently"  # placeholder overwritten below
FAINT   = "6F7789"
GREEN   = "8CD9A0"
AMBER   = "FFC66B"
RED     = "FF7B72"

LOGO = os.path.join(os.environ.get("LOGO_DIR", os.path.dirname(os.path.abspath(__file__))), "logo.png")

FONTS = {
    ("Calibri", False, False): "calibri.ttf",
    ("Calibri", True,  False): "calibrib.ttf",
    ("Calibri", False, True):  "calibrii.ttf",
    ("Calibri", True,  True):  "calibriz.ttf",
    ("Calibri Light", False, False): "calibril.ttf",
    ("Calibri Light", True, False): "calibrib.ttf",
    ("Consolas", False, False): "consola.ttf",
    ("Consolas", True,  False): "consolab.ttf",
    ("Consolas", False, True):  "consolai.ttf",
    ("Consolas", True,  True):  "consolaz.ttf",
}
_fc = {}
def pil_font(name, size_pt, bold=False, italic=False):
    key = (name, bold, italic, size_pt)
    if key not in _fc:
        fn = FONTS.get((name, bold, italic)) or FONTS[("Calibri", bold, italic)]
        _fc[key] = ImageFont.truetype(os.path.join("C:/Windows/Fonts", fn),
                                      max(4, int(round(size_pt / 72 * SCALE))))
    return _fc[key]

def C(hexstr):
    return RGBColor.from_string(hexstr)

def P(hexstr, alpha=255):
    return tuple(int(hexstr[i:i+2], 16) for i in (0, 2, 4)) + (alpha,)

PREVIEW_DIR = None
WARNINGS = []

# ---------------------------------------------------------------- deck class
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(W * EMU_IN))
        self.prs.slide_height = Emu(int(H * EMU_IN))
        self.images = []
        self.sl = None
        self.im = None
        self.dr = None
        self.n = 0

    def slide(self, notes=None):
        self.sl = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        f = self.sl.background.fill
        f.solid(); f.fore_color.rgb = C(BG)
        self.im = Image.new("RGB", (int(W * SCALE), int(H * SCALE)), P(BG)[:3])
        self.dr = ImageDraw.Draw(self.im, "RGBA")
        self.images.append(self.im)
        self.n += 1
        if notes:
            self.sl.notes_slide.notes_text_frame.text = notes
        return self.sl

    # ---- shapes -----------------------------------------------------------
    def _noline(self, sh):
        sh.line.fill.background()

    def _style(self, sh, fill, line, lw):
        if fill is None:
            sh.fill.background()
        else:
            sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
        if line is None:
            self._noline(sh)
        else:
            sh.line.color.rgb = C(line); sh.line.width = Pt(lw)
        sh.shadow.inherit = False

    def rect(self, x, y, w, h, fill=CARD, line=None, lw=1.0):
        sh = self.sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
              Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, line, lw)
        X, Y = x * SCALE, y * SCALE
        self.dr.rectangle([X, Y, X + w * SCALE, Y + h * SCALE],
                          fill=P(fill) if fill else None,
                          outline=P(line) if line else None,
                          width=max(1, int(lw * SCALE / 72)))
        return sh

    def rrect(self, x, y, w, h, fill=CARD, line=None, lw=1.0, rad=0.09):
        sh = self.sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
              Inches(x), Inches(y), Inches(w), Inches(h))
        try:
            sh.adjustments[0] = rad / min(w, h)
        except Exception:
            pass
        self._style(sh, fill, line, lw)
        X, Y = x * SCALE, y * SCALE
        self.dr.rounded_rectangle([X, Y, X + w * SCALE, Y + h * SCALE],
                                  radius=rad * SCALE,
                                  fill=P(fill) if fill else None,
                                  outline=P(line) if line else None,
                                  width=max(1, int(lw * SCALE / 72)))
        return sh

    def oval(self, x, y, w, h, fill=CORAL, line=None, lw=1.0):
        sh = self.sl.shapes.add_shape(MSO_SHAPE.OVAL,
              Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, line, lw)
        X, Y = x * SCALE, y * SCALE
        self.dr.ellipse([X, Y, X + w * SCALE, Y + h * SCALE],
                        fill=P(fill) if fill else None,
                        outline=P(line) if line else None,
                        width=max(1, int(lw * SCALE / 72)))
        return sh

    def hexa(self, x, y, w, h, fill=None, line=CORAL, lw=1.5):
        sh = self.sl.shapes.add_shape(MSO_SHAPE.HEXAGON,
              Inches(x), Inches(y), Inches(w), Inches(h))
        self._style(sh, fill, line, lw)
        # PIL: horizontal hexagon, pointed left/right (adj = 25% of width)
        X, Y, Wp, Hp = x * SCALE, y * SCALE, w * SCALE, h * SCALE
        a = 0.25 * Wp
        pts = [(X + a, Y), (X + Wp - a, Y), (X + Wp, Y + Hp / 2),
               (X + Wp - a, Y + Hp), (X + a, Y + Hp), (X, Y + Hp / 2)]
        self.dr.polygon(pts, fill=P(fill) if fill else None,
                        outline=P(line) if line else None,
                        width=max(1, int(lw * SCALE / 72)))
        return sh

    def cline(self, x1, y1, x2, y2, color=TERMBRD, lw=1.0, dash=None):
        sh = self.sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
              Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        sh.line.color.rgb = C(color); sh.line.width = Pt(lw)
        sh.shadow.inherit = False
        if dash:
            ln = sh.line._get_or_add_ln()
            d = ln.makeelement(qn('a:prstDash'), {'val': dash})
            ln.append(d)
        px = max(1, int(lw * SCALE / 72))
        if dash:
            # crude dash for preview
            import math
            L = math.hypot(x2 - x1, y2 - y1)
            n = max(1, int(L / 0.12))
            for i in range(0, n, 2):
                t0, t1 = i / n, min(1, (i + 1) / n)
                self.dr.line([ (x1 + (x2 - x1) * t0) * SCALE, (y1 + (y2 - y1) * t0) * SCALE,
                               (x1 + (x2 - x1) * t1) * SCALE, (y1 + (y2 - y1) * t1) * SCALE],
                             fill=P(color), width=px)
        else:
            self.dr.line([x1 * SCALE, y1 * SCALE, x2 * SCALE, y2 * SCALE],
                         fill=P(color), width=px)
        return sh

    def pic(self, path, x, y, w=None, h=None):
        kw = {}
        if w: kw["width"] = Inches(w)
        if h: kw["height"] = Inches(h)
        self.sl.shapes.add_picture(path, Inches(x), Inches(y), **kw)
        img = Image.open(path).convert("RGBA")
        iw, ih = img.size
        if w and not h: h = w * ih / iw
        if h and not w: w = h * iw / ih
        img = img.resize((int(w * SCALE), int(h * SCALE)), Image.LANCZOS)
        self.im.paste(img, (int(x * SCALE), int(y * SCALE)), img)

    # ---- text -------------------------------------------------------------
    def text(self, x, y, w, h, content, size=16, font="Calibri", bold=False,
             italic=False, color=TEXT, align="l", valign="t", spacing=1.0,
             space_after=0, wrap=True, shrink_ok=False):
        """content: str | list of paragraphs; paragraph: str | list of runs;
        run: str | (str, dict with size/bold/italic/color/font)."""
        if isinstance(content, str):
            paras = content.split("\n")
        else:
            paras = content
        norm = []
        for p in paras:
            if isinstance(p, str):
                p = [(p, {})]
            else:
                p = [(r, {}) if isinstance(r, str) else r for r in p]
            norm.append(p)

        tb = self.sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = bool(wrap)
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                              "b": MSO_ANCHOR.BOTTOM}[valign]
        amap = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
        for i, runs in enumerate(norm):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = amap[align]
            para.line_spacing = spacing
            if space_after:
                para.space_after = Pt(space_after)
            for t, o in runs:
                r = para.add_run(); r.text = t
                r.font.name = o.get("font", font)
                r.font.size = Pt(o.get("size", size))
                r.font.bold = o.get("bold", bold)
                r.font.italic = o.get("italic", italic)
                r.font.color.rgb = C(o.get("color", color))

        # ---------- PIL rendering with wrap + overflow check ----------
        maxw = w * SCALE * 0.99
        lines = []  # (list of (txt, fnt, colr), line_h_px)
        for runs in norm:
            frags = []
            for t, o in runs:
                fnt = pil_font(o.get("font", font), o.get("size", size),
                               o.get("bold", bold), o.get("italic", italic))
                frags.append((t, fnt, P(o.get("color", color)),
                              o.get("size", size)))
            if not wrap:
                lh = max((f[3] for f in frags), default=size) * 1.22 * spacing
                lines.append(([ (t, f, c) for t, f, c, _ in frags ], lh, False))
                tot = sum(self.dr.textlength(t, font=f) for t, f, c, _ in frags)
                if tot > w * SCALE:
                    WARNINGS.append(f"slide {self.n}: no-wrap line too wide: "
                                    f"{''.join(t for t,_,_,_ in frags)[:60]!r}")
                lines[-1] = (lines[-1][0], lh, True)
                continue
            # word wrap across runs
            words = []
            for t, fnt, colr, sz in frags:
                for wtok in t.split(" "):
                    words.append((wtok, fnt, colr, sz))
            cur, curw = [], 0.0
            spacew = lambda fnt: self.dr.textlength(" ", font=fnt)
            lh = max((f[3] for f in frags), default=size) * 1.22 * spacing
            for wtok, fnt, colr, sz in words:
                ww = self.dr.textlength(wtok, font=fnt)
                add = ww if not cur else ww + spacew(fnt)
                if cur and curw + add > maxw:
                    lines.append((cur, lh, False)); cur, curw = [], 0.0
                    add = ww
                if cur:
                    cur.append((" " + wtok, fnt, colr))
                else:
                    cur.append((wtok, fnt, colr))
                curw += add
            lines.append((cur or [("", frags[0][1] if frags else
                          pil_font(font, size), P(color))], lh, False))
            if space_after:
                lines[-1] = (lines[-1][0], lh + space_after * SCALE / 72 / (SCALE / 72) , True)
                # store marker: extra spacing added below
                lines[-1] = (lines[-1][0], lh, True)
                lines.append(([("", pil_font(font, 4), P(color))],
                              space_after * 1.0, True))
        tot_h = sum(lh / 72 * SCALE for _, lh, _ in lines)
        box_h = h * SCALE
        if tot_h > box_h * 1.02:
            WARNINGS.append(f"slide {self.n}: text overflow ({tot_h/SCALE:.2f}in "
                            f"> {h:.2f}in): {str(norm[0][0][0])[:50]!r}")
        y0 = y * SCALE
        if valign == "m": y0 += (box_h - tot_h) / 2
        if valign == "b": y0 += box_h - tot_h
        for frags, lh, _ in lines:
            lw_px = sum(self.dr.textlength(t, font=f) for t, f, _ in frags)
            if align == "c": xx = x * SCALE + (w * SCALE - lw_px) / 2
            elif align == "r": xx = x * SCALE + w * SCALE - lw_px
            else: xx = x * SCALE
            asc_adj = 0.18 * lh  # rough leading split
            for t, f, colr in frags:
                self.dr.text((xx, y0 + asc_adj / 2), t, font=f, fill=colr)
                xx += self.dr.textlength(t, font=f)
            y0 += lh / 72 * SCALE
        return tb

    def save(self, path):
        self.prs.save(path)
        if PREVIEW_DIR:
            os.makedirs(PREVIEW_DIR, exist_ok=True)
            for i, im in enumerate(self.images, 1):
                im.save(os.path.join(PREVIEW_DIR, f"slide-{i:02d}.png"))

# ---------------------------------------------------------------- helpers
d = Deck()

def kicker(num, label, page=True):
    d.hexa(0.55, 0.48, 0.44, 0.40, fill=None, line=CORAL, lw=1.75)
    d.text(0.55, 0.48, 0.44, 0.40, num, size=13, bold=True, color=CORAL,
           align="c", valign="m", font="Consolas", wrap=False)
    d.text(1.13, 0.48, 8.0, 0.40, label.upper(), size=12.5, bold=True,
           color=MUTED, valign="m", wrap=False)
    if page:
        d.text(12.45, 7.08, 0.5, 0.3, f"{d.n:02d}", size=10, color=FAINT,
               align="r", font="Consolas", wrap=False)

def title(t, y=0.98, size=30, w=12.2):
    d.text(0.55, y, w, 0.62, t, size=size, bold=True, color=TEXT, wrap=False)

def term_card(x, y, w, h, rows, title_txt=None):
    """rows: list of (kind, text): 'cmd','out','mut','com','amber'"""
    d.rrect(x, y, w, h, fill=TERM, line=TERMBRD, lw=1.0, rad=0.10)
    for i, c in enumerate(("FF6E67", "FFC66B", "8CD9A0")):
        d.oval(x + 0.18 + i * 0.17, y + 0.15, 0.09, 0.09, fill=c)
    if title_txt:
        d.text(x + 0.75, y + 0.10, w - 0.9, 0.22, title_txt, size=10.5,
               color=FAINT, font="Consolas", wrap=False)
    yy = y + 0.44
    for kind, txt in rows:
        if kind == "gap":
            yy += 0.10; continue
        col = {"cmd": TEXT, "out": GREEN, "mut": MUTED, "com": FAINT,
               "amber": AMBER}[kind]
        runs = []
        if kind == "cmd":
            runs = [("$ ", {"color": CORAL, "bold": True}), (txt, {})]
        else:
            runs = [(txt, {})]
        d.text(x + 0.25, yy, w - 0.45, 0.26, [runs], size=12, font="Consolas",
               color=col, wrap=False)
        yy += 0.255
    return yy

def bullet_rows(x, y, w, items, gap=None, tsize=15.5, dsize=12.5, dot=True):
    """items: list of (head, desc). Returns bottom y."""
    yy = y
    for head, desc in items:
        if dot:
            d.oval(x, yy + 0.075, 0.10, 0.10, fill=CORAL)
        d.text(x + (0.28 if dot else 0), yy - 0.02, w - 0.28, 0.30, head,
               size=tsize, bold=True, color=TEXT, wrap=False)
        if desc:
            d.text(x + (0.28 if dot else 0), yy + 0.27, w - 0.28, 0.55, desc,
                   size=dsize, color=MUTED, spacing=1.02)
            yy += (gap or 0.78)
        else:
            yy += (gap or 0.45)
    return yy

# ================================================================ SLIDE 1
d.slide(notes=(
    "Setup before people arrive: terminal open, ollama serve running, models "
    "pre-pulled. This title slide stays up while people come in.\n"
    "Talk: 30-45 min + Q&A, live demos on the 16GB laptop."))
d.pic(LOGO, 0.95, 1.55, h=3.5)
x0 = 5.0
d.text(x0, 1.95, 7.6, 0.35, "CODE FOR BIELEFELD · JULI 2026", size=13,
       bold=True, color=CORAL, wrap=False)
d.text(x0, 2.32, 7.9, 1.15, "Open Source LLMs", size=54, bold=True,
       color=TEXT, wrap=False)
d.text(x0, 3.42, 7.8, 0.85,
       "Capable AI on your own laptop — no cloud, no API key, no data leaving the room.",
       size=19, color=MUTED, spacing=1.1)
d.text(x0, 4.42, 7.8, 0.35, "Live demos throughout · Ollama · a normal 16 GB laptop",
       size=13, color=FAINT, wrap=False)
term_card(x0, 5.15, 7.35, 1.30, [
    ("cmd", 'ollama run gemma4:e2b'),
    ("mut", '>>> Erkläre in 3 Sätzen, warum Bielefeld existiert.'),
    ("out", 'Bielefeld existiert – und das lässt sich sogar lokal beweisen. _'),
])

# ================================================================ SLIDE 2
d.slide(notes=(
    "Open with the live demo BEFORE any slides — run gemma4:e2b, ask it "
    "something fun in German (Gemma 4 supports 140+ languages).\n"
    "Land the point hard: this is running on THIS laptop. WiFi can be off.\n"
    "Fallback: screen recording of this exact demo."))
kicker("00", "Hook · live demo first")
title("This just ran on this laptop")
term_card(0.55, 1.85, 7.35, 3.5, [
    ("cmd", 'ollama run gemma4:e2b'),
    ("mut", '>>> Erkläre in 3 Sätzen, warum Bielefeld existiert.'),
    ("gap", ""),
    ("out", 'Bielefeld existiert seit 1214 als Kaufmannsstadt am'),
    ("out", 'Übergang durch den Teutoburger Wald. Die Sparrenburg'),
    ("out", 'bewachte den Pass, die Leinenweber brachten den Handel.'),
    ("out", 'Dass die Stadt nicht existiere, ist nur ein Internet-Witz.'),
    ("gap", ""),
    ("com", '# gemma4:e2b · 2.3B params · ~2 GB · CPU only'),
    ("gap", ""),
    ("com", '# ~40 tokens/s on this CPU — fast enough to read along'),
], title_txt="live — no network needed")
d.text(0.55, 5.52, 7.35, 0.4,
       "Flugmodus an — die Demo läuft trotzdem.", size=13, italic=True,
       color=FAINT, wrap=False)
bx = 8.35
for i, (head, sub) in enumerate([
        ("Runs entirely on this laptop", "2 GB download · 16 GB consumer laptop · no GPU"),
        ("No cloud, no API key", "No account, no rate limits, no monthly bill"),
        ("No data leaves the room", "Privacy by construction — DSGVO-friendly")]):
    y = 1.85 + i * 1.35
    d.rrect(bx, y, 4.45, 1.15, fill=CARD, rad=0.10)
    d.text(bx + 0.25, y + 0.18, 4.0, 0.3, head, size=15.5, bold=True,
           color=CORAL if i == 2 else TEXT, wrap=False)
    d.text(bx + 0.25, y + 0.52, 4.0, 0.5, sub, size=12, color=MUTED, spacing=1.02)

# ================================================================ SLIDE 3
d.slide(notes=(
    "TODO before talk: double-check current Gemma license terms vs Apache 2.0.\n"
    "Point for this audience: 'open weights' is the norm — you can run and "
    "fine-tune, but training data/recipe usually stay closed.\n"
    "Civic-tech angle: DSGVO + Verwaltung is the strongest why-care here."))
kicker("01", "What “open” actually means")
title("“Open source LLM” usually means open weights")
rows = [("Open weights", "You get the trained parameters — download, run, fine-tune. This is what “open” means for almost every model.", "the norm", GREEN),
        ("Open data", "The training data is published. Very few models (e.g. OLMo-style research projects).", "rare", AMBER),
        ("Open training code", "The full recipe to reproduce the model from scratch.", "rare", AMBER)]
for i, (head, desc, tag, tc) in enumerate(rows):
    y = 1.80 + i * 1.06
    d.rrect(0.55, y, 7.35, 0.92, fill=CARD, rad=0.09)
    d.text(0.85, y + 0.13, 2.6, 0.3, head, size=16, bold=True, wrap=False)
    d.text(0.85, y + 0.47, 6.2, 0.4, desc, size=11.5, color=MUTED, spacing=1.0)
    d.rrect(6.55, y + 0.16, 1.05, 0.34, fill=TERM, rad=0.17)
    d.text(6.55, y + 0.16, 1.05, 0.34, tag, size=10.5, bold=True, color=tc,
           align="c", valign="m", wrap=False)
d.text(0.55, 5.12, 7.35, 0.75, [
    [("Licenses still differ:  ", {"bold": True, "color": TEXT}),
     ("Apache 2.0 (Qwen — do anything)  ·  Gemma license  ·  Llama license (usage terms attached)",
      {"color": MUTED})]], size=12.5, spacing=1.05)
d.text(8.35, 1.72, 4.5, 0.3, "WHY CARE?", size=12.5, bold=True, color=CORAL,
       wrap=False)
bullet_rows(8.35, 2.15, 4.45, [
    ("Privacy — DSGVO!", "Citizen data never leaves your machine"),
    ("Cost & offline", "No per-token bill, works without internet"),
    ("No vendor lock-in", "The model file is yours, forever"),
    ("Fine-tunable", "Adapt it to your domain (→ LoRA, later)")], gap=0.80)
d.text(8.35, 5.45, 4.45, 0.6,
       "A very good fit for civic tech and public administration.",
       size=12.5, italic=True, color=FAINT, spacing=1.05)

# ================================================================ SLIDE 4
d.slide(notes=(
    "TODO: pick current DeepSeek / Llama / Mistral versions to mention (one "
    "line each).\n"
    "Media idea: show the Ollama model library or HF leaderboard live if WiFi "
    "works — have a screenshot as fallback."))
kicker("02", "The landscape · mid-2026")
title("Who’s who in open models")
# Gemma card
d.rrect(0.55, 1.80, 6.0, 3.55, fill=CARD, rad=0.12)
d.text(0.90, 2.05, 5.3, 0.4, "Gemma 4", size=22, bold=True, wrap=False)
d.text(0.90, 2.50, 5.3, 0.3, "Google — “the multimodal generalist”",
       size=12.5, color=CORAL, wrap=False)
bullet_rows(0.90, 2.98, 5.4, [
    ("E2B · E4B · 12B Unified · 26B MoE · 31B", "sizes for phone → workstation"),
    ("Multimodal", "text + image (12B Unified also audio) · 256K context"),
    ("140+ languages", "notably strong German")], gap=0.72, tsize=13.5, dsize=11.5)
# Qwen card
d.rrect(6.80, 1.80, 6.0, 3.55, fill=CARD, rad=0.12)
d.text(7.15, 2.05, 5.3, 0.4, "Qwen 3.6", size=22, bold=True, wrap=False)
d.text(7.15, 2.50, 5.3, 0.3, "Alibaba — “the coding monster”",
       size=12.5, color=CORAL, wrap=False)
bullet_rows(7.15, 2.98, 5.4, [
    ("27B dense", "one size, tuned hard for code & agents"),
    ("77.2% SWE-bench Verified", "on par with frontier closed models"),
    ("Apache 2.0", "the most permissive license there is")], gap=0.72,
    tsize=13.5, dsize=11.5)
d.rrect(0.55, 5.60, 12.25, 0.95, fill=CARD2, rad=0.10)
d.text(0.90, 5.72, 11.7, 0.7, [
    [("Also on the map:  ", {"bold": True}),
     ("DeepSeek (open reasoning models)  ·  Llama (Meta’s ecosystem)  ·  "
      "Mistral (European, efficient)", {"color": MUTED})]],
    size=13.5, valign="m", spacing=1.0)

# ================================================================ SLIDE 5
d.slide(notes=(
    "THE story of the talk: capability got compressed into small models.\n"
    "Killer stat: Qwen3.6-27B beats its own 397B predecessor on coding "
    "benchmarks — small + recent > big + old."))
kicker("03", "The real story")
title("From datacenter to laptop in six years")
tl_y = 2.75
d.cline(0.9, tl_y, 12.4, tl_y, color=TERMBRD, lw=2)
pts = [("2020", "GPT-3", "175B · closed", "datacenter only"),
       ("2023", "Llama 2", "70B · open-ish", "workstation"),
       ("2025", "Gemma 3 / Qwen 3", "4–32B · open", "laptop"),
       ("2026", "Gemma 4 E2B/E4B", "2–4B · open · multimodal", "any laptop, even phones")]
for i, (yr, name, spec, where) in enumerate(pts):
    x = 1.35 + i * 3.35
    d.oval(x - 0.09, tl_y - 0.09, 0.18, 0.18,
           fill=CORAL if i == 3 else BG, line=CORAL, lw=1.75)
    d.text(x - 1.3, tl_y - 0.75, 2.6, 0.3, yr, size=14, bold=True,
           color=CORAL, align="c", wrap=False)
    d.text(x - 1.45, tl_y + 0.28, 2.9, 0.3, name, size=15, bold=True,
           align="c", wrap=False)
    d.text(x - 1.45, tl_y + 0.62, 2.9, 0.3, spec, size=11.5, color=MUTED,
           align="c", wrap=False)
    d.text(x - 1.45, tl_y + 0.90, 2.9, 0.3, where, size=11.5, color=FAINT,
           align="c", italic=True, wrap=False)
d.rrect(1.7, 4.55, 9.95, 1.85, fill=CARD, line=CORALDK, lw=1.25, rad=0.14)
d.text(2.1, 4.85, 9.2, 0.6,
       "Qwen 3.6-27B beats its own 397B predecessor on coding benchmarks",
       size=20, bold=True, color=TEXT, spacing=1.05)
d.text(2.1, 5.68, 9.2, 0.4, [
    [("small + recent  >  big + old", {"color": CORAL, "bold": True,
                                       "font": "Consolas", "size": 15}),
     ("   — capability per parameter keeps climbing", {"color": MUTED,
                                                       "size": 13})]],
    wrap=False)

# ================================================================ SLIDE 6
d.slide(notes=(
    "DEMO BLOCK 1 (~8 min). Pre-pull all models; close browsers so 12B fits; "
    "have screen recordings as fallback for every demo.\n"
    "Run the same prompt at 3 sizes with --verbose, watch tokens/sec drop and "
    "quality rise. Second terminal: ollama ps shows RAM footprint + 100% CPU.\n"
    "Then the wall: qwen3.6:27b Q4 is ~16GB weights alone → doesn't fit → "
    "perfect segue to quantization & hardware.\n"
    "Multimodal teaser: image of the Rathaus into gemma4:e4b."))
kicker("04", "Live demo · small models on this laptop")
title("Same prompt, three sizes — watch tokens/sec")
term_card(0.55, 1.80, 7.35, 2.65, [
    ("cmd", 'ollama run gemma4:e2b --verbose "Schreibe eine'),
    ("mut", '    Python-Funktion: Ist ein Jahr ein Schaltjahr?"'),
    ("com", '# … eval rate: ~XX tokens/s   ← measure live'),
    ("gap", ""),
    ("cmd", 'ollama run gemma4:e4b --verbose "…"'),
    ("cmd", 'ollama run gemma4:12b --verbose "…"'),
    ("com", '# speed drops, quality rises'),
], title_txt="demo 1 — size vs speed")
term_card(0.55, 4.72, 7.35, 1.93, [
    ("cmd", 'ollama run gemma4:e4b "Was ist auf diesem'),
    ("mut", '    Bild zu sehen?" ./bielefeld-rathaus.jpg'),
    ("com", '# multimodal teaser — image in, answer out'),
    ("com", '# gemma4 reads forms, timetables, handwriting'),
], title_txt="demo 2 — vision")
term_card(8.35, 1.80, 4.45, 1.95, [
    ("cmd", 'ollama ps'),
    ("mut", 'NAME         SIZE   PROC'),
    ("out", 'gemma4:12b   8.1GB  100% CPU'),
    ("com", '# RAM footprint, live'),
], title_txt="second terminal")
d.rrect(8.35, 4.00, 4.45, 2.65, fill=CARD, line=None, rad=0.12)
d.text(8.65, 4.22, 3.9, 0.35, "The wall", size=17, bold=True, color=AMBER,
       wrap=False)
d.text(8.65, 4.62, 3.95, 0.62, [
    [("$ ", {"color": CORAL, "bold": True, "font": "Consolas", "size": 12}),
     ("ollama pull qwen3.6:27b", {"font": "Consolas", "size": 12})]],
    wrap=False)
d.text(8.65, 4.98, 3.9, 1.5,
       "Q4 weights alone ≈ 16 GB — more than this laptop’s entire RAM.\n"
       "Doesn’t fit. Why? And what would?", size=12.5, color=MUTED,
       spacing=1.1, space_after=4)

# ================================================================ SLIDE 7
d.slide(notes=(
    "TODO: run test_token_speeds.py beforehand and fill in the real tokens/s "
    "numbers + your quality ratings (build the table from YOUR measurements — "
    "more credible than benchmark screenshots). Show the script briefly.\n"
    "TODO: decide the big-hardware story — rent a GPU for real Qwen 3.6 "
    "numbers, or use published benchmarks.\n"
    "Verify gemma4:12b actually runs OK in 16GB (close browsers!)."))
kicker("05", "The numbers")
title("Speed vs. RAM vs. quality — measured here")
tab_x, tab_w = 0.55, 12.25
cols = [3.1, 1.3, 1.2, 1.5, 2.6, 2.55]
heads = ["MODEL", "PARAMS", "QUANT", "RAM", "TOKENS/S (this laptop)", "QUALITY"]
rows7 = [
    ("gemma4:e2b",  "2.3B", "Q4", "~3 GB",  "→ measured live", "fine for chat"),
    ("gemma4:e4b",  "4B",   "Q4", "~5 GB",  "→ measured live", "good all-rounder"),
    ("gemma4:12b",  "12B",  "Q4", "~8 GB",  "→ measured live", "noticeably better"),
    ("qwen3.6:27b", "27B",  "Q4", "~17 GB", "doesn’t fit",   "—"),
    ("qwen3.6:27b  (rented GPU)", "27B", "Q8", "30–60 GB", "30–60 t/s",
     "frontier coding"),
]
yy = 1.80
xx = tab_x + 0.25
for j, htxt in enumerate(heads):
    d.text(xx, yy, cols[j] - 0.1, 0.3, htxt, size=10.5, bold=True,
           color=FAINT, wrap=False)
    xx += cols[j]
yy += 0.38
for i, r in enumerate(rows7):
    if i % 2 == 0:
        d.rrect(tab_x, yy - 0.065, tab_w, 0.52, fill=CARD2, rad=0.06)
    xx = tab_x + 0.25
    for j, val in enumerate(r):
        colr = TEXT
        fnt = "Calibri"; bd = False; sz = 13.5
        if j == 0: fnt = "Consolas"; sz = 12.5
        if j == 4: colr = CORAL if val.startswith("→") else (RED if val.startswith("doesn") else GREEN); bd = True
        if j == 5: colr = MUTED
        d.text(xx, yy, cols[j] - 0.1, 0.34, val, size=sz, bold=bd, color=colr,
               font=fnt, wrap=False)
        xx += cols[j]
    yy += 0.52
d.rrect(0.55, 4.75, 12.25, 1.85, fill=CARD, rad=0.12)
d.text(0.90, 4.98, 11.6, 0.35, "Rule of thumb", size=15, bold=True,
       color=CORAL, wrap=False)
d.text(0.90, 5.38, 11.6, 0.45,
       [[("RAM  ≈  parameters × bytes per weight  +  context overhead",
          {"font": "Consolas", "size": 16, "bold": True})]], wrap=False)
d.text(0.90, 5.92, 11.6, 0.45,
       "FP16 = 2 bytes/param · Q4 ≈ 0.5 bytes/param   →   a 27B model: 54 GB in FP16, ~15 GB in Q4",
       size=13, color=MUTED, wrap=False)

# ================================================================ SLIDE 8
d.slide(notes=(
    "The one formula that explains the whole hardware market. LLM generation "
    "is memory-bandwidth-bound, not compute-bound: every generated token "
    "reads all weights once.\n"
    "CPU dual-channel DDR5 ~100GB/s → 4B Q4 (~2.5GB) at ~40 t/s is pleasant; "
    "27B Q4 (~16GB) at ~6 t/s is painful.\n"
    "Optional demo: ollama ps shows 100% CPU here; screenshot of the same on "
    "a GPU box showing 100% GPU + the t/s difference."))
kicker("06", "Hardware")
title("One formula explains the entire market")
d.rrect(0.55, 1.75, 12.25, 1.15, fill=TERM, line=CORALDK, lw=1.25, rad=0.12)
d.text(0.55, 1.75, 12.25, 1.15,
       [[("tokens/s   ≈   memory bandwidth  ÷  model size in bytes",
          {"font": "Consolas", "size": 22, "bold": True, "color": CORAL})]],
       align="c", valign="m", wrap=False)
d.text(0.55, 3.05, 12.25, 0.35,
       "Generating one token reads every weight once — inference is memory-bound, not compute-bound.",
       size=13.5, color=MUTED, align="c", wrap=False)
bars = [("Laptop CPU · DDR5", 100, MUTED), ("Ryzen AI Max 395 · unified", 256, MUTED),
        ("Mac M4 Max · unified", 546, MUTED), ("RTX 3090 · GDDR6X", 936, CORAL),
        ("RTX 5090 · GDDR7", 1792, CORAL)]
by = 3.70; bw_max = 7.6; vmax = 1792.0
for name, val, colr in bars:
    d.text(0.55, by, 3.15, 0.3, name, size=12.5, color=TEXT, align="r",
           wrap=False)
    bw = max(0.25, bw_max * val / vmax)
    d.rrect(3.90, by + 0.015, bw, 0.26, fill=colr, rad=0.06)
    d.text(3.98 + bw, by, 1.4, 0.3, f"{val} GB/s", size=12, bold=True,
           color=TEXT, font="Consolas", wrap=False)
    by += 0.47
d.text(3.90, by + 0.10, 8.8, 0.4, [
    [("On this laptop:  ", {"bold": True, "size": 12.5}),
     ("4B Q4 ≈ 40 t/s — pleasant      27B Q4 ≈ 6 t/s — painful",
      {"font": "Consolas", "size": 12.5, "color": MUTED})]], wrap=False)

# ================================================================ SLIDE 9
d.slide(notes=(
    "Prices are mid-2026 street prices, approximate. TODO: refresh shortly "
    "before the talk — GPU prices move fast.\n"
    "TODO: check rent-vs-buy break-even (at ~€0.50/h a used 3090 pays for "
    "itself after roughly 1,300 hours).\n"
    "Talking points: used RTX 3090 = the community Geheimtipp (~€1,200 total "
    "build). The VRAM cliff: if the model doesn't fit in VRAM it spills to "
    "system RAM and speed falls off a cliff. And: CPU is fine for small "
    "models — that's the real revolution, not the €3,000 GPU."))
kicker("07", "Hardware · what does speed cost?")
title("The buyer’s table (mid-2026)")
cols9 = [3.4, 1.9, 1.6, 1.7, 2.4, 1.25]
heads9 = ["OPTION", "MEMORY", "BANDWIDTH", "PRICE", "SWEET SPOT", "QWEN 27B?"]
rows9 = [
    ("This laptop (CPU)", "16 GB shared", "100 GB/s", "owned", "≤ 4B models", ("no", RED)),
    ("Used RTX 3090", "24 GB VRAM", "936 GB/s", "€600–700", "best t/s per €", ("yes · fast", GREEN)),
    ("RTX 5090", "32 GB VRAM", "1792 GB/s", "€2,600+", "70B Q4, 45+ t/s", ("yes · v. fast", GREEN)),
    ("Ryzen AI Max+ 395", "128 GB unified", "256 GB/s", "€1,500–2,000", "big models, slowly", ("yes · medium", AMBER)),
    ("Mac Studio", "512 GB unified", "546–819 GB/s", "€2,500+", "70B @ 20–30 t/s", ("yes · fast", GREEN)),
    ("Rented GPU (vast.ai)", "any", "any", "€0.30–0.70/h", "try before buying", ("yes", GREEN)),
]
yy = 1.78
xx = 0.80
for j, htxt in enumerate(heads9):
    d.text(xx, yy, cols9[j] - 0.1, 0.3, htxt, size=10.5, bold=True,
           color=FAINT, wrap=False)
    xx += cols9[j]
yy += 0.36
for i, r in enumerate(rows9):
    if i % 2 == 0:
        d.rrect(0.55, yy - 0.06, 12.25, 0.50, fill=CARD2, rad=0.06)
    xx = 0.80
    for j, val in enumerate(r):
        if j == 5:
            t, colr = val
            d.text(xx, yy, cols9[j] - 0.05, 0.34, t, size=12.5, bold=True,
                   color=colr, wrap=False)
        else:
            bd = (j == 0)
            d.text(xx, yy, cols9[j] - 0.1, 0.34, val, size=12.5, bold=bd,
                   color=TEXT if j == 0 else MUTED, wrap=False)
        xx += cols9[j]
    yy += 0.50
chips = [("Used 3090 = the Geheimtipp", "same 24 GB as a 4090 at ⅓ the price"),
         ("Beware the VRAM cliff", "doesn’t fit in VRAM → speed falls off a cliff"),
         ("CPU is fine ≤ 4B", "that’s the revolution — not the €3,000 GPU")]
for i, (head, sub) in enumerate(chips):
    x = 0.55 + i * 4.20
    d.rrect(x, 5.55, 3.95, 1.15, fill=CARD, rad=0.10)
    d.text(x + 0.22, 5.70, 3.55, 0.3, head, size=13.5, bold=True, color=CORAL,
           wrap=False)
    d.text(x + 0.22, 6.03, 3.55, 0.55, sub, size=11.5, color=MUTED, spacing=1.0)

# ================================================================ SLIDE 10
d.slide(notes=(
    "Intuition first: weights are just numbers; quantization stores them "
    "sloppier. JPEG analogy carries the whole section.\n"
    "Live if time: same model in q4 and q8 → compare speed, RAM, and (usually) "
    "barely different output. An over-quantized Q2 model giving worse answers "
    "is more memorable than any chart.\n"
    "Mention: GGUF file format, llama.cpp is the engine underneath Ollama.\n"
    "Media idea: 3Blue1Brown clips for 'weights are numbers' intuition; "
    "perplexity-vs-quant chart from llama.cpp docs."))
kicker("08", "Quantization")
title("JPEG for neural networks")
d.text(0.55, 1.75, 5.6, 1.6,
       "Weights are just numbers. Quantization stores them with fewer bits — "
       "lossy, like JPEG, but you barely notice.",
       size=16, color=TEXT, spacing=1.15)
d.text(0.55, 3.0, 5.6, 0.4,
       [[("FP16 → Q4 = 4× smaller, ~same answers",
          {"font": "Consolas", "size": 13.5, "color": CORAL, "bold": True})]],
       wrap=False)
qb = [("FP16", 8.0, 100, "reference", MUTED),
      ("Q8", 4.3, 100, "indistinguishable", GREEN),
      ("Q4", 2.4, 98, "sweet spot ← use this", CORAL),
      ("Q2", 1.5, 84, "visibly degraded", RED)]
qy = 1.85
d.text(6.75, qy - 0.32, 3.0, 0.28, "SIZE (4B MODEL)", size=10, bold=True,
       color=FAINT, wrap=False)
d.text(10.6, qy - 0.32, 2.2, 0.28, "QUALITY", size=10, bold=True, color=FAINT,
       wrap=False)
for name, gb, qual, verdict, colr in qb:
    d.text(6.10, qy, 0.62, 0.3, name, size=13, bold=True, font="Consolas",
           color=TEXT, align="r", wrap=False)
    bw = 2.5 * gb / 8.0
    d.rrect(6.90, qy + 0.02, bw, 0.24, fill=TERMBRD, rad=0.05)
    d.text(6.98 + bw, qy, 1.0, 0.28, f"{gb:.1f} GB", size=11, color=MUTED,
           font="Consolas", wrap=False)
    qw = 2.0 * qual / 100.0
    d.rrect(10.60, qy + 0.02, qw, 0.24, fill=colr, rad=0.05)
    d.text(6.90, qy + 0.30, 5.9, 0.26, verdict, size=10.5, italic=True,
           color=colr, wrap=False)
    qy += 0.72
term_card(0.55, 4.95, 7.35, 1.75, [
    ("cmd", 'ollama pull gemma4:e2b-it-q8_0    # vs default q4_K_M'),
    ("com", '# same prompt on both → compare RAM, speed, output'),
    ("com", '# spoiler: you will not spot the difference'),
], title_txt="demo — q8 vs q4")
d.rrect(8.35, 4.95, 4.45, 1.75, fill=CARD2, rad=0.10)
d.text(8.65, 5.18, 3.9, 0.3, "Under the hood", size=13.5, bold=True,
       color=TEXT, wrap=False)
d.text(8.65, 5.52, 3.95, 1.0,
       "GGUF = the quantized model file format.\nllama.cpp = the engine "
       "beneath Ollama.", size=12, color=MUTED, spacing=1.12)

# ================================================================ SLIDE 11
d.slide(notes=(
    "The core math slide. Training needs weights + gradients + Adam optimizer "
    "states + activations ≈ 10–20× the memory of quantized inference — and it "
    "can't be aggressively quantized because tiny gradient updates vanish in "
    "4-bit.\n"
    "Compute side: Gemma-class pretraining = thousands of GPUs × weeks.\n"
    "TODO: fun back-of-envelope — 'training Gemma 4 on this laptop ≈ X years' "
    "(FLOPs of pretraining ÷ laptop FLOPs).\n"
    "Media idea: datacenter photo vs this laptop, side by side."))
kicker("09", "Inference vs. training")
title("You can run it. You cannot train it.")
d.text(0.55, 1.70, 6.4, 0.35, "MEMORY FOR A 4B MODEL", size=11, bold=True,
       color=FAINT, wrap=False)
iy = 2.15
d.text(0.55, iy, 2.1, 0.3, "Inference Q4", size=13, bold=True, wrap=False)
d.rrect(2.85, iy + 0.02, 0.32, 0.26, fill=GREEN, rad=0.05)
d.text(3.28, iy, 2.6, 0.3, "2.5 GB — easy", size=12.5, color=GREEN,
       font="Consolas", wrap=False)
iy += 0.52
d.text(0.55, iy, 2.1, 0.3, "Inference FP16", size=13, bold=True, wrap=False)
d.rrect(2.85, iy + 0.02, 1.0, 0.26, fill=MUTED, rad=0.05)
d.text(3.97, iy, 2.6, 0.3, "8 GB — tight", size=12.5, color=MUTED,
       font="Consolas", wrap=False)
iy += 0.52
d.text(0.55, iy, 2.1, 0.3, "Full training", size=13, bold=True, color=AMBER,
       wrap=False)
segs = [("weights", 8, "C4655F"), ("gradients", 8, "FF867F"),
        ("optimizer", 16, "FFC66B"), ("activations", 18, "FFE3A3")]
sx = 2.85
for name, gb, colr in segs:
    swd = gb * 0.0625
    d.rect(sx, iy + 0.02, swd, 0.26, fill=colr)
    sx += swd
d.text(sx + 0.12, iy, 2.6, 0.3, "50+ GB — no", size=12.5, bold=True,
       color=RED, font="Consolas", wrap=False)
iy += 0.42
lx = 2.85
lfont = pil_font("Calibri", 10.5)
for name, gb, colr in segs:
    d.rrect(lx, iy + 0.045, 0.13, 0.13, fill=colr, rad=0.03)
    lbl = f"{name} {gb}"
    d.text(lx + 0.19, iy, 1.6, 0.24, lbl, size=10.5, color=MUTED, wrap=False)
    lx += 0.19 + d.dr.textlength(lbl, font=lfont) / SCALE + 0.22
iy += 0.36
# 16GB laptop line
lap_x = 2.85 + 16 * 0.0625
d.cline(lap_x, 2.05, lap_x, 4.05, color=FAINT, lw=1.25, dash="dash")
d.text(lap_x + 0.08, 3.93, 2.5, 0.26, "16 GB — this laptop", size=10,
       color=FAINT, italic=True, wrap=False)
d.text(0.55, 4.42, 6.6, 0.35, "Training a 27B: ~350+ GB → multi-GPU cluster.",
       size=13, color=MUTED, wrap=False)
d.rrect(0.55, 4.95, 6.6, 1.7, fill=CARD, line=CORALDK, lw=1.0, rad=0.12)
d.text(0.55, 4.95, 6.6, 1.7, [
    [("Inference is reading the book.", {"size": 17, "bold": True})],
    [("Training is writing it.", {"size": 17, "bold": True, "color": CORAL})]],
    align="c", valign="m", spacing=1.25, wrap=False)
d.text(7.60, 1.70, 5.2, 0.35, "WHY THE GAP?", size=11, bold=True, color=FAINT,
       wrap=False)
for i, (head, sub) in enumerate([
        ("10–20× the memory of Q4 inference",
         "backprop keeps gradients + Adam states + activations alongside the weights"),
        ("Training can’t be 4-bit",
         "tiny gradient updates vanish in coarse number formats — precision is required"),
        ("And it’s compute, too",
         "Gemma-class pretraining: thousands of GPUs for weeks — this laptop: centuries")]):
    y = 2.15 + i * 1.5
    d.rrect(7.60, y, 5.20, 1.3, fill=CARD2, rad=0.10)
    d.text(7.85, y + 0.16, 4.75, 0.55, head, size=14.5, bold=True, spacing=1.0)
    d.text(7.85, y + 0.62, 4.75, 0.6, sub, size=11.5, color=MUTED, spacing=1.05)

# ================================================================ SLIDE 12
d.slide(notes=(
    "Bridge: '…but there IS a trick to customize models at home.'\n"
    "LoRA: freeze the base weights, train two tiny low-rank matrices on the "
    "side (~0.1–1% of params). QLoRA = same on a quantized base → fine-tuning "
    "Gemma 4 E4B fits in ~8–12 GB / free Colab T4.\n"
    "Demo plan: train the LoRA BEFORE the talk (don't train live — too slow, "
    "too fragile), show the Unsloth notebook briefly, demo the result in "
    "Ollama.\n"
    "TODO: decide the fun dataset + build it (~100–500 examples is enough for "
    "a style LoRA): Bielefelder Platt / Amtsdeutsch / CfB style."))
kicker("10", "The loophole")
title("Fine-tuning at home: LoRA")
# diagram
d.rrect(0.55, 2.0, 3.6, 2.5, fill=CARD, rad=0.12)
d.text(0.55, 2.25, 3.6, 0.4, "Base model", size=16, bold=True, align="c",
       wrap=False)
d.text(0.55, 2.65, 3.6, 0.3, "4B params · quantized", size=11.5, color=MUTED,
       align="c", wrap=False)
d.rrect(1.25, 3.10, 2.2, 0.55, fill=TERM, rad=0.10)
d.text(1.25, 3.10, 2.2, 0.55, "FROZEN", size=13, bold=True, color=MUTED,
       align="c", valign="m", font="Consolas", wrap=False)
d.text(0.55, 3.85, 3.6, 0.3, "0% trained", size=11, color=FAINT, align="c",
       italic=True, wrap=False)
d.text(4.30, 2.9, 0.6, 0.5, "+", size=30, bold=True, color=CORAL, align="c",
       wrap=False)
d.rrect(5.00, 2.35, 1.15, 0.85, fill=CORAL, rad=0.10)
d.text(5.00, 2.35, 1.15, 0.85, "A", size=20, bold=True, color=BG, align="c",
       valign="m", wrap=False)
d.rrect(5.00, 3.35, 1.15, 0.85, fill=CORAL, rad=0.10)
d.text(5.00, 3.35, 1.15, 0.85, "B", size=20, bold=True, color=BG, align="c",
       valign="m", wrap=False)
d.text(4.72, 4.35, 1.75, 0.6, "two tiny matrices\n~0.1–1% of params",
       size=10.5, color=MUTED, align="c", spacing=1.05)
d.text(6.25, 2.9, 0.6, 0.5, "=", size=30, bold=True, color=CORAL, align="c",
       wrap=False)
d.rrect(6.95, 2.55, 2.35, 1.45, fill=CARD, line=CORAL, lw=1.5, rad=0.12)
d.text(6.95, 2.75, 2.35, 0.4, "Your model", size=15, bold=True, align="c",
       wrap=False)
d.text(6.95, 3.15, 2.35, 0.6, "custom style,\nsame brain", size=11,
       color=MUTED, align="c", spacing=1.05)
d.text(0.55, 5.02, 8.75, 0.55, [
    [("QLoRA:", {"bold": True, "color": CORAL}),
     ("  train on top of the quantized base  →  Gemma 4 E4B fine-tunes in "
      "~8–12 GB — a free Colab T4 is enough (Unsloth).", {"color": MUTED})]],
    size=13.5, spacing=1.05)
d.rrect(0.55, 5.65, 8.75, 1.05, fill=CARD2, rad=0.10)
d.text(0.85, 5.78, 8.2, 0.8, [
    [("Demo LoRA idea:  ", {"bold": True, "size": 13}),
     ("Gemma antwortet auf Bielefelder Platt — oder in astreinem Amtsdeutsch. "
      "Trained before the talk, shown live in Ollama.",
      {"color": MUTED, "size": 12.5})]], valign="m", spacing=1.05)
d.rrect(9.65, 2.0, 3.15, 2.1, fill=CARD, rad=0.12)
d.text(9.90, 2.18, 2.7, 0.3, "LoRA can", size=14, bold=True, color=GREEN,
       wrap=False)
d.text(9.90, 2.55, 2.7, 1.4, "style & tone\noutput format\ndomain vocabulary",
       size=12.5, color=TEXT, spacing=1.25)
d.rrect(9.65, 4.35, 3.15, 1.75, fill=CARD, rad=0.12)
d.text(9.90, 4.53, 2.7, 0.3, "LoRA can’t", size=14, bold=True, color=RED,
       wrap=False)
d.text(9.90, 4.90, 2.7, 1.1, "new knowledge at scale\nnew capabilities",
       size=12.5, color=TEXT, spacing=1.25)

# ================================================================ SLIDE 13
d.slide(notes=(
    "Message: pick the model for the task, not the biggest one.\n"
    "Gemma demo: photo of a paper form / Fahrplan / handwritten note → "
    "extract structured JSON, fully local → 'DSGVO-clean form scanning' is "
    "the civic-tech hook.\n"
    "Qwen demo (recorded from bigger hardware or via API): a real coding task "
    "both models attempt — Gemma E4B struggles, Qwen nails it.\n"
    "TODO: verify Qwen 3.6 27B is really text-only (some sources list MMMU "
    "scores → suggests vision?).\n"
    "Optional if time: Qwen plugged into Cline/aider = 'local GitHub "
    "Copilot'."))
kicker("11", "Capabilities")
title("Pick the model for the task — not the biggest one")
d.rrect(0.55, 1.80, 6.0, 3.6, fill=CARD, rad=0.12)
d.text(0.90, 2.02, 5.3, 0.35, "Gemma 4 — sees, hears, speaks German",
       size=16, bold=True, wrap=False)
d.text(0.90, 2.42, 5.3, 0.3, "the multimodal generalist", size=12,
       color=CORAL, italic=True, wrap=False)
bullet_rows(0.90, 2.90, 5.4, [
    ("Photo of a paper form → JSON", "scan forms into structured data, fully local"),
    ("140+ languages, strong German", "12B Unified adds audio input"),
    ("DSGVO-clean by design", "the civic-tech / Verwaltung sweet spot")],
    gap=0.78, tsize=13.5, dsize=11.5)
d.rrect(6.80, 1.80, 6.0, 3.6, fill=CARD, rad=0.12)
d.text(7.15, 2.02, 5.3, 0.35, "Qwen 3.6 — writes serious code",
       size=16, bold=True, wrap=False)
d.text(7.15, 2.42, 5.3, 0.3, "the specialist coder", size=12, color=CORAL,
       italic=True, wrap=False)
bullet_rows(7.15, 2.90, 5.4, [
    ("77.2% SWE-bench Verified", "real repository bugs, frontier-level"),
    ("Same task, side by side", "Gemma E4B struggles — Qwen nails it (demo)"),
    ("Local GitHub Copilot", "plug into Cline or aider — private code stays private")],
    gap=0.78, tsize=13.5, dsize=11.5)
d.rrect(0.55, 5.65, 12.25, 1.0, fill=CARD2, rad=0.10)
d.text(0.90, 5.78, 11.7, 0.75, [
    [("Honorable mentions:  ", {"bold": True}),
     ("Whisper (speech→text, runs anywhere)  ·  small embedding models "
      "(local RAG)  ·  DeepSeek-R1-style reasoning", {"color": MUTED})]],
    size=13.5, valign="m", spacing=1.0)

# ================================================================ SLIDE 14
d.slide(notes=(
    "Wrap-up, ~2 min. Land the three takeaways, then the invitation: "
    "ollama.com, one command, running in five minutes.\n"
    "Close with the Code for Bielefeld project idea: local document "
    "processing or RAG over the Ratsinformationssystem — collect hands/"
    "interest for a follow-up evening."))
kicker("12", "Wrap-up")
title("Three things to take home")
takes = [
    ("1", "Capable LLMs run on normal laptops — today",
     "Quantization (Q4) makes it possible; small recent models beat big old ones"),
    ("2", "Inference is cheap. Training is astronomical.",
     "LoRA / QLoRA is the middle path — customize a model on a free Colab"),
    ("3", "Model choice = task fit",
     "Gemma 4 for multimodal & German · Qwen 3.6 for code · size vs. speed")]
for i, (n, head, sub) in enumerate(takes):
    y = 1.85 + i * 1.18
    d.rrect(0.55, y, 7.9, 1.02, fill=CARD, rad=0.10)
    d.hexa(0.85, y + 0.27, 0.52, 0.48, fill=None, line=CORAL, lw=1.75)
    d.text(0.85, y + 0.27, 0.52, 0.48, n, size=16, bold=True, color=CORAL,
           align="c", valign="m", font="Consolas", wrap=False)
    d.text(1.60, y + 0.15, 6.6, 0.32, head, size=15.5, bold=True, wrap=False)
    d.text(1.60, y + 0.52, 6.6, 0.4, sub, size=11.5, color=MUTED, wrap=False)
term_card(0.55, 5.50, 7.9, 1.2, [
    ("cmd", 'ollama run gemma4:e4b        # ollama.com'),
    ("com", '# one command → running in 5 minutes'),
], title_txt="try it yourself — tonight")
d.rrect(8.85, 1.85, 3.95, 4.85, fill=CARD, line=CORALDK, lw=1.0, rad=0.14)
d.text(9.15, 2.10, 3.4, 0.6, "A project for\nCode for Bielefeld?", size=17,
       bold=True, spacing=1.05)
bullet_rows(9.15, 3.05, 3.5, [
    ("Local document AI", "paper forms → structured data"),
    ("RAG over the Ratsinformationssystem", "ask the city council’s documents"),
    ("All DSGVO-clean", "runs on hardware we own")], gap=0.85, tsize=12.5,
    dsize=11)
d.text(9.15, 6.05, 3.4, 0.5, "Interested? Let’s talk after.", size=12.5,
       italic=True, color=CORAL, wrap=False)

# ================================================================ SLIDE 15
d.slide(notes=(
    "Q&A. Sources listed for the follow-up email / repo.\n"
    "Have the terminal still open — best Q&A answers are live demos."))
d.pic(LOGO, 5.62, 0.85, h=2.3)
d.text(0.55, 3.35, 12.25, 0.7, "Danke! — Questions?", size=38, bold=True,
       align="c", wrap=False)
d.text(0.55, 4.15, 12.25, 0.35,
       "The terminal is still open — ask, and we’ll try it live.",
       size=14, color=CORAL, align="c", italic=True, wrap=False)
srcs = [
    "Google — Gemma 4 announcement (blog.google)",
    "Unsloth — Gemma 4 fine-tuning docs (unsloth.ai)",
    "Ollama model library (ollama.com/library)",
    "The Decoder — Qwen 3.6-27B vs its 397B predecessor",
    "llm-stats.com — Qwen 3.6 27B benchmarks",
    "corelab.tech — LLM GPU buyer's guide (Apr 2026)",
    "modemguides.com — Ryzen AI Max 395 reality check",
    "pinggy.io — hardware for self-hosting LLMs, 2026",
]
d.text(2.3, 4.95, 4.4, 0.3, "SOURCES", size=10, bold=True, color=FAINT,
       wrap=False)
for i, s in enumerate(srcs):
    xcol = 2.3 if i < 4 else 7.1
    yrow = 5.28 + (i % 4) * 0.33
    d.text(xcol, yrow, 4.6, 0.28, s, size=10.5, color=MUTED, wrap=False)

# ---------------------------------------------------------------- save
if __name__ == "__main__":
    if "--preview" in sys.argv:
        PREVIEW_DIR = sys.argv[sys.argv.index("--preview") + 1]
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "OpenSource-LLMs-CodeForBielefeld.pptx")
    d.save(out)
    print("written:", out)
    if WARNINGS:
        print("\n-- layout warnings --")
        for wmsg in WARNINGS:
            print(" !", wmsg)
    else:
        print("no layout warnings")
